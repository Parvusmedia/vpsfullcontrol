#!/usr/bin/env python3
"""Generate editable Fly456 PowerPoint for Kima outreach.

Positioning: named flight monitoring (origin, destination, dates, max price)
+ Telegram alerts to book fast. Economy and Business cabin products.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(0x0F, 0x1E, 0x2D)
TEAL = RGBColor(0x00, 0x6E, 0x82)
ACCENT = RGBColor(0xE6, 0x78, 0x28)
LIGHT = RGBColor(0xF6, 0xF8, 0xFA)
MUTED = RGBColor(0x5F, 0x6E, 0x78)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x19, 0x23, 0x2D)

OUTS = [
    Path(__file__).resolve().parent / "Fly456_Kima_OnePager.pptx",
    Path(__file__).resolve().parents[1] / "Fly456_Kima_OnePager.pptx",
    Path("/opt/cursor/artifacts/Fly456_Kima_OnePager.pptx"),
]


def add_bg(slide, prs, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp_tree = slide.shapes._spTree
    sp = shape._element
    sp_tree.remove(sp)
    sp_tree.insert(2, sp)


def bar_left(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), Inches(7.5))
    s.fill.solid()
    s.fill.fore_color.rgb = TEAL
    s.line.fill.background()


def top_strip(slide):
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.22), 0, Inches(13.333 - 0.22), Inches(0.45)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = LIGHT
    s.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.08), Inches(10), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = "Fly456  ·  Pre-seed note for Kima Ventures"
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED
    p.font.name = "Calibri"


def add_text(slide, l, t, w, h, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    return tb


def add_bullets(slide, l, t, w, h, lines, size=16):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + line
        p.font.size = Pt(size)
        p.font.color.rgb = DARK
        p.font.name = "Calibri"
        p.space_after = Pt(10)


def card(slide, l, t, w, h, title, lines):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = LIGHT
    s.line.fill.background()
    a = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(0.06))
    a.fill.solid()
    a.fill.fore_color.rgb = TEAL
    a.line.fill.background()
    add_text(slide, l + 0.2, t + 0.2, w - 0.4, 0.35, title, 16, True, NAVY)
    tb = slide.shapes.add_textbox(Inches(l + 0.2), Inches(t + 0.6), Inches(w - 0.4), Inches(h - 0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = DARK
        p.font.name = "Calibri"
        p.space_after = Pt(4)


def content_header(slide, prs, eyebrow, title, page):
    add_bg(slide, prs, WHITE)
    bar_left(slide)
    top_strip(slide)
    add_text(slide, 12.2, 0.08, 0.9, 0.3, f"{page}/6", 11, False, MUTED, PP_ALIGN.RIGHT)
    add_text(slide, 0.5, 0.7, 12, 0.3, eyebrow.upper(), 12, True, TEAL)
    add_text(slide, 0.5, 1.05, 12, 0.7, title, 28, True, NAVY)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 Cover
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, prs, NAVY)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.28), prs.slide_height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    accent = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.28), Inches(6.55), prs.slide_width - Inches(0.28), Inches(0.14)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()
    add_text(
        s,
        0.7,
        1.2,
        11,
        0.4,
        "PRE-SEED  ·  FLIGHT MONITORING  ·  TELEGRAM ALERTS",
        14,
        False,
        RGBColor(0xB4, 0xC8, 0xD2),
    )
    add_text(s, 0.7, 1.8, 11, 0.9, "Fly456", 54, True, WHITE)
    add_text(
        s,
        0.7,
        2.8,
        11,
        1.2,
        "You set origin, destination, dates and max price.\nWe monitor flights and notify you on Telegram — so you can book fast.",
        20,
        False,
        RGBColor(0xC8, 0xDC, 0xE1),
    )
    add_text(
        s,
        0.7,
        4.3,
        11,
        0.4,
        "Two products: Economy cabin  ·  Business cabin",
        16,
        True,
        ACCENT,
    )
    add_text(
        s,
        0.7,
        5.0,
        11,
        0.4,
        "Asking: ~€150k pre-seed  ·  fly456.com  ·  @fly456bot",
        15,
        True,
        RGBColor(0xA0, 0xB4, 0xBE),
    )
    add_text(
        s,
        0.7,
        5.6,
        11,
        0.4,
        "Emiliano Tichauer  ·  Founder",
        14,
        False,
        RGBColor(0xA0, 0xB4, 0xBE),
    )
    add_text(s, 0.7, 6.85, 2, 0.3, "1/6", 12, False, RGBColor(0x8C, 0xA0, 0xAA))

    # 2 Problem
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(
        s, prs, "Problem", "Good fares disappear before travelers can act", 2
    )
    add_bullets(
        s,
        0.6,
        2.0,
        12,
        4.5,
        [
            "Flight prices move in hours. Checking sites manually is slow and easy to miss.",
            "Generic deal feeds are noisy — they are not your route, dates or budget.",
            "By the time you see a fare, seats and price are often already gone.",
        ],
        18,
    )

    # 3 Solution
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(
        s, prs, "Solution", "Named monitoring + Telegram alert to book fast", 3
    )
    card(
        s,
        0.5,
        2.0,
        4.0,
        4.2,
        "You define the flight",
        [
            "Origin",
            "Destination",
            "Dates",
            "Maximum price",
            "",
            "Fly456 watches that exact brief.",
        ],
    )
    card(
        s,
        4.75,
        2.0,
        4.0,
        4.2,
        "We notify on Telegram",
        [
            "Alert lands in your Telegram inbox",
            "via @fly456bot",
            "Built to help you reserve fast",
            "before the fare moves again",
            "",
            "We do not sell tickets.",
        ],
    )
    card(
        s,
        9.0,
        2.0,
        3.8,
        4.2,
        "Two cabin products",
        [
            "Economy — monitor economy",
            "cabin fares for your route",
            "",
            "Business — monitor business",
            "cabin fares for your route",
            "",
            "Same alert logic, different cabin.",
        ],
    )

    # 4 Product
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(s, prs, "Product", "Economy and Business — same monitoring, different cabin", 4)
    card(
        s,
        0.5,
        2.0,
        6.0,
        4.2,
        "Fly456 Economy",
        [
            "Monitors economy-cabin opportunities",
            "Alert = origin + destination + dates + max price",
            "Telegram notification when a match appears",
            "For leisure and price-sensitive travelers",
            "who need speed, not another search tab",
        ],
    )
    card(
        s,
        6.8,
        2.0,
        6.0,
        4.2,
        "Fly456 Business",
        [
            "Monitors business-cabin opportunities",
            "Same named criteria, business fares",
            "Telegram notification to book fast",
            "For travelers who want premium cabin",
            "when the price finally drops into range",
        ],
    )

    # 5 Traction
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(s, prs, "Traction", "Just launched July 2026, no marketing budget", 5)
    add_bullets(
        s,
        0.6,
        2.0,
        12,
        3.8,
        [
            "Currently 6 free channels live on Telegram for niche audiences (LatAm expats living in EU and US).",
            "Currently few paying Premium subscribers.",
            "Plans: Economy 5€/month | 20€/year · Business 12€/month | 99€/year. Stripe checkout in Telegram.",
            "Product live: fly456.com + @fly456bot",
        ],
        17,
    )
    card(
        s,
        0.5,
        5.5,
        12.3,
        1.5,
        "The ask (~€150k)",
        [
            "Sharper monitoring · faster Telegram alerts · grow paid Economy + Business · keep burn low",
        ],
    )

    # 6 Team
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_header(s, prs, "Team", "Operator-builder shipping the full loop", 6)
    add_bullets(
        s,
        0.6,
        2.0,
        12,
        4.0,
        [
            "Emiliano Tichauer — Founder",
            "Built Parvus Media since 2013 (AdTech, data, performance products)",
            "Owns the loop end-to-end: flight monitoring, Telegram UX, Stripe billing",
            "Contact: [email]  ·  LinkedIn: linkedin.com/in/etichauer  ·  fly456.com",
        ],
        18,
    )

    for out in OUTS:
        out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out))
        print(f"Wrote {out}")


if __name__ == "__main__":
    build()
