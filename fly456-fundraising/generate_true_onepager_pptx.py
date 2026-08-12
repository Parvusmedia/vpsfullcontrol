#!/usr/bin/env python3
"""True 1-slide Fly456 one-pager for Kima."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x0F, 0x1E, 0x2D)
TEAL = RGBColor(0x00, 0x6E, 0x82)
ACCENT = RGBColor(0xE6, 0x78, 0x28)
LIGHT = RGBColor(0xF6, 0xF8, 0xFA)
MUTED = RGBColor(0x5F, 0x6E, 0x78)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x19, 0x23, 0x2D)
OUTS = [
    Path(__file__).resolve().parent / "Fly456_Kima_OnePager.pptx",
    Path(__file__).resolve().parents[1] / "Fly456_Kima_OnePager.pptx",
    Path("/opt/cursor/artifacts/Fly456_Kima_OnePager.pptx"),
]

def build():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = WHITE; bg.line.fill.background()
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.18), prs.slide_height)
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
    hdr = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.18), 0, prs.slide_width - Inches(0.18), Inches(1.35))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = NAVY; hdr.line.fill.background()
    acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.18), Inches(1.35), prs.slide_width - Inches(0.18), Inches(0.06))
    acc.fill.solid(); acc.fill.fore_color.rgb = ACCENT; acc.line.fill.background()

    def tb(l, t, w, h, text, size=12, bold=False, color=DARK, align=PP_ALIGN.LEFT):
        box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame; tf.word_wrap = True
        for i, line in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line; p.font.size = Pt(size); p.font.bold = bold
            p.font.color.rgb = color; p.font.name = "Calibri"; p.alignment = align; p.space_after = Pt(2)

    def card(l, t, w, h, title, body_lines):
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = LIGHT; sh.line.fill.background()
        top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(0.05))
        top.fill.solid(); top.fill.fore_color.rgb = TEAL; top.line.fill.background()
        tb(l + 0.15, t + 0.12, w - 0.3, 0.28, title, 13, True, TEAL)
        box = s.shapes.add_textbox(Inches(l + 0.15), Inches(t + 0.42), Inches(w - 0.3), Inches(h - 0.55))
        tf = box.text_frame; tf.word_wrap = True
        for i, line in enumerate(body_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line; p.font.size = Pt(12); p.font.color.rgb = DARK; p.font.name = "Calibri"; p.space_after = Pt(3)

    tb(0.45, 0.18, 8, 0.3, "FLY456  ·  ONE-PAGER FOR KIMA  ·  PRE-SEED", 11, False, RGBColor(0xB4, 0xC8, 0xD2))
    tb(0.45, 0.48, 9, 0.45, "Named flight monitoring. Telegram alert. Book fast.", 26, True, WHITE)
    tb(0.45, 0.95, 10, 0.3, "Origin · Destination · Dates · Max price  →  alert on @fly456bot", 14, False, RGBColor(0xC8, 0xDC, 0xE1))
    tb(10.2, 0.35, 2.8, 0.7, "Ask\n~€150k", 20, True, ACCENT, PP_ALIGN.RIGHT)

    card(0.4, 1.6, 6.2, 2.15, "PROBLEM", [
        "• Fares move in hours — manual checking is too slow.",
        "• Generic deal feeds are noisy (not your brief).",
        "• When you finally see it, the fare is often gone.",
    ])
    card(6.85, 1.6, 6.1, 2.15, "PRODUCT", [
        "• You set: origin, destination, dates, max price.",
        "• Fly456 monitors that exact brief.",
        "• Telegram notifies you via @fly456bot to reserve fast.",
        "• We don’t sell tickets — you book with airline/OTA.",
    ])
    card(0.4, 3.9, 6.2, 1.55, "FLY456 ECONOMY", [
        "Economy-cabin monitoring for your named brief.",
        "Pricing: 5€/month  ·  20€/year",
    ])
    card(6.85, 3.9, 6.1, 1.55, "FLY456 BUSINESS", [
        "Business-cabin monitoring for your named brief.",
        "Pricing: 12€/month  ·  99€/year",
    ])
    card(0.4, 5.6, 6.2, 1.6, "TRACTION (JUL 2026 LAUNCH · NO ADS)", [
        "• 6 free Telegram channels (LatAm expats in EU/US).",
        "• Early Premium — few paying subscribers so far.",
        "• Live: fly456.com + @fly456bot (Stripe in Telegram).",
    ])
    card(6.85, 5.6, 6.1, 1.6, "USE OF FUNDS (~€150k)", [
        "• BizDev hire  ·  paid acquisition tests",
        "• Free → Premium conversion  ·  affiliate pilots",
        "• Prove LTV/CAC with low burn",
    ])
    foot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.18), Inches(7.22), prs.slide_width - Inches(0.18), Inches(0.28))
    foot.fill.solid(); foot.fill.fore_color.rgb = NAVY; foot.line.fill.background()
    tb(0.45, 7.24, 12.5, 0.25, "Emiliano Tichauer — Founder  ·  emiliano@tichauer.es  ·  linkedin.com/in/etichauer  ·  fly456.com", 11, False, WHITE)

    for out in OUTS:
        out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out))
        print(f"Wrote {out}")

if __name__ == "__main__":
    build()
